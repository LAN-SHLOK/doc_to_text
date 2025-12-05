# core/doc_to_text.py
"""
Robust doc_to_text engine that:
- Uses PyMuPDF (fitz) when available for fast text extraction
- Falls back to pdfminer.six for embedded-text PDFs when PyMuPDF is unavailable
- Uses pytesseract for OCR and makes EasyOCR optional (fallback to pytesseract if missing)
"""
import os
import sys
import re
from pathlib import Path
from typing import List, Tuple, Dict
import tempfile

# Optional imports with graceful fallbacks
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except Exception:
    fitz = None
    HAS_PYMUPDF = False

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
    HAS_PDFMINER = True
except Exception:
    pdfminer_extract_text = None
    HAS_PDFMINER = False

# OCR libs
from PIL import Image
import pytesseract

# EasyOCR is optional (heavy dependency). We'll try import but not require it.
try:
    import easyocr
    HAS_EASYOCR = True
except Exception:
    easyocr = None
    HAS_EASYOCR = False

# image processing
try:
    import cv2
    import numpy as np
    HAS_CV2 = True
except Exception:
    cv2 = None
    np = None
    HAS_CV2 = False

# other libs
try:
    import pandas as pd
except Exception:
    pd = None

# Spell checker optional
try:
    from spellchecker import SpellChecker
    GLOBAL_SPELLER = SpellChecker()
except Exception:
    GLOBAL_SPELLER = None

# config
EASYOCR_LANGS = ['en']
DO_SPELL_CORRECTION = False


# Lazy initializer for easyocr.Reader
_reader_easyocr = None
def get_easyocr_reader():
    global _reader_easyocr
    if not HAS_EASYOCR:
        return None
    if _reader_easyocr is None:
        _reader_easyocr = easyocr.Reader(EASYOCR_LANGS, gpu=False)
    return _reader_easyocr


# ---------- Utilities ----------
def preprocess_image_for_ocr(pil_image: Image.Image) -> Image.Image:
    """
    If cv2 available, use denoising + adaptive threshold. Else return original image.
    """
    if not HAS_CV2:
        return pil_image.convert('RGB')
    img = np.array(pil_image.convert('RGB'))
    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    denoised = cv2.fastNlMeansDenoising(gray, h=10)
    th = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                               cv2.THRESH_BINARY, 15, 10)
    return Image.fromarray(th)


def ocr_pytesseract_with_conf(img: Image.Image) -> List[Tuple[str, float]]:
    """
    Use pytesseract image_to_data (if pandas available) to extract per-word confidences.
    If pandas isn't available, return words with 0.0 confidence.
    """
    try:
        tsv = pytesseract.image_to_data(img, output_type=pytesseract.Output.DATAFRAME)
    except Exception:
        raw = pytesseract.image_to_string(img)
        words = [w for w in re.split(r'\s+', raw.strip()) if w]
        return [(w, 0.0) for w in words]

    words_conf = []
    if tsv is None or (hasattr(tsv, 'empty') and tsv.empty):
        return words_conf

    for _, row in tsv.iterrows():
        txt = str(row.get('text', '')).strip()
        conf_val = row.get('conf', None)
        conf = -1.0
        try:
            if pd is not None and pd.isna(conf_val):
                conf = -1.0
            else:
                conf = float(conf_val)
        except Exception:
            try:
                conf = float(str(conf_val))
            except Exception:
                conf = -1.0
        if txt and conf != -1.0:
            words_conf.append((txt, max(0.0, min(1.0, conf / 100.0))))
    return words_conf


def ocr_easyocr_with_conf(img: Image.Image) -> List[Tuple[str, float]]:
    """
    Use EasyOCR if available. If not, return empty list to let pytesseract drive OCR.
    """
    if not HAS_EASYOCR:
        return []
    reader = get_easyocr_reader()
    if reader is None:
        return []
    arr = np.array(img.convert('RGB'))
    results = reader.readtext(arr)
    out = []
    for bbox, text, conf in results:
        parts = re.split(r'\s+', text.strip())
        for p in parts:
            if p:
                try:
                    c = float(conf)
                    if c > 1.0:
                        c = c / 100.0
                except Exception:
                    c = 0.0
                out.append((p, max(0.0, min(1.0, c))))
    return out


def merge_word_lists_by_conf(wlist1, wlist2):
    merged = []
    confs = []
    n = max(len(wlist1), len(wlist2))
    for i in range(n):
        w1, c1 = (("", 0.0) if i >= len(wlist1) else wlist1[i])
        w2, c2 = (("", 0.0) if i >= len(wlist2) else wlist2[i])
        if c1 >= c2:
            chosen, chosen_conf = w1, c1
        else:
            chosen, chosen_conf = w2, c2
        if not chosen:
            if w1:
                chosen, chosen_conf = w1, c1
            elif w2:
                chosen, chosen_conf = w2, c2
        if chosen:
            merged.append(chosen)
            confs.append(float(chosen_conf))
    avg_conf = float(np.mean(confs)) if HAS_CV2 and confs else (float(sum(confs))/len(confs) if confs else 0.0)
    return merged, avg_conf


def words_to_text(words):
    return " ".join(words)


# ---------- Format extractors ----------
def extract_text_from_docx(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception:
        return ""


def extract_text_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        out_lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        out_lines.append(txt)
        return "\n".join(out_lines)
    except Exception:
        return ""


def extract_text_from_text_pdf(path: Path) -> str:
    """
    Use PyMuPDF when available, else pdfminer.six fallback.
    """
    if HAS_PYMUPDF:
        out = []
        doc = fitz.open(str(path))
        for page in doc:
            ptxt = page.get_text("text")
            if ptxt and ptxt.strip():
                out.append(ptxt.strip())
        return "\n\n".join(out)
    if HAS_PDFMINER:
        try:
            return pdfminer_extract_text(str(path)) or ""
        except Exception:
            return ""
    return ""


def pdf_to_images(path: Path, dpi:int=300):
    from pdf2image import convert_from_path
    return convert_from_path(str(path), dpi=dpi)


def extract_text_by_ocr_from_image(img: Image.Image) -> Tuple[str, float]:
    prep = preprocess_image_for_ocr(img)
    try:
        t_words = ocr_pytesseract_with_conf(prep)
    except Exception:
        t_words = []
    try:
        e_words = ocr_easyocr_with_conf(prep)
    except Exception:
        e_words = []

    # fallback logic
    if not t_words and not e_words:
        raw = pytesseract.image_to_string(prep)
        words = [w for w in re.split(r'\s+', raw.strip()) if w]
        return words_to_text(words), 0.0
    if not t_words:
        merged_words = [w for w,_ in e_words]
        avg_conf = float(sum([c for _,c in e_words]) / len(e_words)) if e_words else 0.0
        return words_to_text(merged_words), avg_conf
    if not e_words:
        merged_words = [w for w,_ in t_words]
        avg_conf = float(sum([c for _,c in t_words]) / len(t_words)) if t_words else 0.0
        return words_to_text(merged_words), avg_conf

    merged_words, avg_conf = merge_word_lists_by_conf(t_words, e_words)
    return words_to_text(merged_words), avg_conf


def extract_text_from_pdf_with_mixed_strategy(path: Path) -> Tuple[str, float]:
    text = extract_text_from_text_pdf(path)
    if len(text.strip()) > 50:
        return text, 0.99
    # scanned -> OCR per page
    try:
        images = pdf_to_images(path, dpi=300)
    except Exception:
        images = []
    all_texts = []
    confs = []
    for img in images:
        txt, conf = extract_text_by_ocr_from_image(img)
        if txt.strip():
            all_texts.append(txt)
            confs.append(conf)
    combined = "\n\n".join(all_texts)
    avg_conf = float(sum(confs)/len(confs)) if confs else 0.0
    return combined, avg_conf


def extract_text_from_image_file(path: Path) -> Tuple[str, float]:
    img = Image.open(str(path))
    return extract_text_by_ocr_from_image(img)


# ---------- Important details ----------
def extract_important_details(text: str) -> Dict[str, List[str]]:
    details = {"emails": [], "phones": [], "key_values": [], "dates": []}
    details['emails'] = list(set(re.findall(r'[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}', text)))
    details['phones'] = list(set(re.findall(r'\+?\d[\d\-\s]{6,}\d', text)))
    kvs=[]
    for line in text.splitlines():
        if ':' in line:
            left,right = line.split(':',1)
            left = left.strip(); right = right.strip()
            if left and right and len(left) < 120:
                kvs.append({left: right})
    details['key_values'] = kvs
    dates = re.findall(r'\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}-\d{2}-\d{2}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s*\d{0,4})\b', text, flags=re.IGNORECASE)
    details['dates'] = list(set(dates))
    return details


# ---------- Main orchestration ----------
def convert_any_to_text(filepath: str, do_spell_correct: bool = DO_SPELL_CORRECTION) -> Tuple[str, float, str]:
    p = Path(filepath)
    if not p.exists():
        raise FileNotFoundError(filepath)

    suffix = p.suffix.lower()
    combined_text = ""
    estimated_conf = 0.0

    if suffix in ['.docx']:
        combined_text = extract_text_from_docx(p)
        estimated_conf = 0.999 if combined_text else 0.0

    elif suffix in ['.doc']:
        raise RuntimeError(".doc (binary) not supported. Convert to .docx first.")

    elif suffix in ['.pptx']:
        combined_text = extract_text_from_pptx(p)
        estimated_conf = 0.995 if combined_text else 0.0

    elif suffix in ['.pdf']:
        combined_text, estimated_conf = extract_text_from_pdf_with_mixed_strategy(p)

    elif suffix in ['.jpg','jpeg','png','tiff','bmp','webp']:
        combined_text, estimated_conf = extract_text_from_image_file(p)

    else:
        try:
            with open(p, 'r', encoding='utf-8') as f:
                combined_text = f.read()
            estimated_conf = 0.999
        except Exception:
            combined_text, estimated_conf = extract_text_from_image_file(p)

    if do_spell_correct and GLOBAL_SPELLER and combined_text.strip():
        # light correction
        tokens = re.split(r'(\W+)', combined_text)
        corrected=[]
        for tok in tokens:
            if re.match(r'^[A-Za-z]{3,}$', tok):
                corr = GLOBAL_SPELLER.correction(tok)
                corrected.append(corr if corr else tok)
            else:
                corrected.append(tok)
        combined_text = "".join(corrected)

    out_path = p.with_suffix('.txt')
    with open(out_path, 'w', encoding='utf-8') as fout:
        fout.write(combined_text)

    return str(out_path), float(estimated_conf), combined_text
